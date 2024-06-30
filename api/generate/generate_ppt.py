from flask import render_template, stream_with_context, current_app
from elasticsearch_client import get_elasticsearch_chat_message_history
from pptx import Presentation
from pptx.util import Pt
from llm_integrations import get_llm
from langchain.prompts import ChatPromptTemplate, MessagesPlaceholder
from value_map import INDEX_CHAT_HISTORY, DONE_TAG, CHAT_DONE_TAG, SEARCH_DONE_TAG, TOTAL_TAG, FACET_TAG, SOURCE_TAG, SESSION_ID_TAG, MSG_SOURCE_TAG, MSG_STATS_TAG, EXPORT_PATH, FILEID_TAG
from langchain.schema import StrOutputParser
import time
import json
from uuid import uuid4
import os
from pptx.util import Inches
import pandas as pd
from io import StringIO
import mistune

def estimate_font_size(text):
    # Estimate the amount of text and return a suitable font size
    # This is a very simple estimation, you may need to adjust it based on your needs
        if len(text) < 100:
            return Pt(24)
        elif len(text) < 200:
            return Pt(18)
        elif len(text) < 500:
            return Pt(14)
        else:
            return Pt(10)

def add_table_to_slide(slide, df):
    # Define table dimensions
    rows, cols = df.shape
    left = Inches(1.0)
    top = Inches(1.5)
    width = Inches(6.0)
    height = Inches(0.8)

    # Add table to slide
    table = slide.shapes.add_table(rows+1, cols, left, top, width, height).table

    # Add the header
    for i, column in enumerate(df.columns):
        cell = table.cell(0, i)
        cell.text = column

    # Add the rest of the data frame
    for i in range(rows):
        for j in range(cols):
            cell = table.cell(i+1, j)
            cell.text = str(df.iloc[i, j])

def is_dataframe(content_text):
    try:
        df = pd.read_table(StringIO(content_text), sep="|")
        df = df.dropna(how="all", axis=0)  # Drop rows
        df = df.dropna(how="all", axis=1)  # Drop columns
        df.fillna("")
        return True
    except pd.errors.ParserError:
        return False

@stream_with_context
def generate_ppt(session_id, start_time=time.time()):
    yield f"data: {SESSION_ID_TAG} {session_id}\n\n"
    current_app.logger.debug("Chat session ID: %s", session_id)

    chat_history = get_elasticsearch_chat_message_history(
        INDEX_CHAT_HISTORY, session_id
    )

    # template_presentation = Presentation("api\\templates\\template.pptx")
    # slide_layouts = template_presentation.slide_layouts

    presentation = Presentation("api\\templates\\blank_template.pptx")

    slide_layout = presentation.slide_layouts[0]  # 0 is title slide
    slide = presentation.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = chat_history.messages[0].content

    for i in range(0, len(chat_history.messages), 2):
        # Skip if there's not a pair of messages
        if i+1 >= len(chat_history.messages):
            break

        # Add a slide with a title and content layout
        slide_layout = presentation.slide_layouts[1]  # 1 is title and content
        slide = presentation.slides.add_slide(slide_layout)

        # Use the human's message as the title and the AI's message as the content
        title = slide.shapes.title
        title.text = chat_history.messages[i].content if chat_history.messages[i].type == "human" else chat_history.messages[i+1].content

        # Set the font size of the title
        for paragraph in title.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = estimate_font_size(title.text)

        content = slide.placeholders[1]  # placeholders[1] is the content placeholder
        content_text = chat_history.messages[i].content if chat_history.messages[i].type == "ai" else chat_history.messages[i+1].content

        # Check if the content is a dataframe in markdown format
        # if is_dataframe(content_text):
        #     df = pd.read_table(StringIO(content_text), sep="|")
        #     add_table_to_slide(slide, df)
        if content_text.endswith('.png') and os.path.exists(EXPORT_PATH + "/" + content_text):
            # Add the image to the slide
            slide.shapes.add_picture(EXPORT_PATH + "/" + content_text, Inches(1), Inches(1), width=Inches(5.5), height=Inches(4.5))
        else:
            # Convert markdown to HTML if it's markdown
            # if content_text.startswith('#') or '**' in content_text or '*' in content_text or '`' in content_text:
            #     # Convert markdown to plain text
            #     markdown_renderer = mistune.Markdown()
            #     content_text = markdown_renderer(content_text)
            #     # Convert the list of dictionaries to a string
            #     content_text = json.dumps(content_text)
            content.text = content_text

        # Set the font size of the content
        for paragraph in content.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = estimate_font_size(content.text)

    # Save the presentation
    file_id = chat_history.messages[0].content+"_"+str(uuid4())+".pptx"
    presentation.save(EXPORT_PATH+"/"+file_id)

    token = {"token":"Here is your presentation."}
    yield f"data: {json.dumps(token)}\n\n"

    fileid_token = {"file_id":file_id}

    yield f"data: {FILEID_TAG} {json.dumps(fileid_token)}\n\n"

    ttft = round((time.time() - start_time)/60,2)
    ttf = round((time.time() - start_time)/60,2)
    yield f'data: {MSG_STATS_TAG} {json.dumps({"ttf":ttf, "ttft":ttft})}\n\n'

    yield f"data: {CHAT_DONE_TAG}\n\n"
    yield f"data: {DONE_TAG}\n\n"